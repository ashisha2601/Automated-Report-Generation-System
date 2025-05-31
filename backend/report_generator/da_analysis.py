import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os
from datetime import datetime
from typing import Dict, List, Optional
import logging
from pydantic import BaseModel
import uuid

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class DailyAssessmentAnalyzer:
    def __init__(self, file_path: str):
        """Initialize the analyzer with the path to the DA data file."""
        self.file_path = file_path
        self.data = None
        self.load_data()

    def load_data(self):
        """Load data from Excel/CSV file."""
        try:
            if self.file_path.endswith('.csv'):
                self.data = pd.read_csv(self.file_path)
            else:
                self.data = pd.read_excel(self.file_path)
            logger.info(f"Data loaded successfully from {self.file_path}")
        except Exception as e:
            logger.error(f"Error loading data: {str(e)}")
            raise

    def apply_filters(self, filters: Dict[str, List[str]]) -> None:
        """Apply filters to the dataset."""
        if not filters:
            return

        for column, values in filters.items():
            if column in self.data.columns and values:
                self.data = self.data[self.data[column].isin(values)]
                logger.info(f"Applied filter: {column} = {values}")

    def generate_summary_statistics(self) -> Dict:
        """Generate summary statistics for the dataset."""
        try:
            summary = {
                'total_records': len(self.data),
                'date_range': {
                    'start': self.data['date'].min().strftime('%Y-%m-%d'),
                    'end': self.data['date'].max().strftime('%Y-%m-%d')
                },
                'metrics': {}
            
                }

            # Calculate metrics for numeric columns
            numeric_columns = self.data.select_dtypes(include=[np.number]).columns
            for col in numeric_columns:
                summary['metrics'][col] = {
                    'mean': self.data[col].mean(),
                    'median': self.data[col].median(),
                    'std': self.data[col].std(),
                    'min': self.data[col].min(),
                    'max': self.data[col].max()
                }

            return summary
        except Exception as e:
            logger.error(f"Error generating summary statistics: {str(e)}")
            raise

    def _create_pre_post_comparison_plot(self, output_dir: str) -> Optional[str]:
        """Create Pre/Post score comparison plot by category and save it."""
        try:
            # Ensure required columns exist (adjust names if necessary)
            required_cols = ['Category', 'AssessmentType', 'Score']
            if not all(col in self.data.columns for col in required_cols):
                logger.warning(f"Skipping Pre/Post plot: Missing one or more required columns: {required_cols}")
                return None
                
            # Filter for 'Pre' and 'Post' and group data
            pre_post_data = self.data[self.data['AssessmentType'].isin(['Pre', 'Post'])].copy()
            if pre_post_data.empty:
                logger.warning("Skipping Pre/Post plot: No 'Pre' or 'Post' data found in 'AssessmentType' column.")
                return None

            # Calculate average scores by Category and AssessmentType
            avg_scores = pre_post_data.groupby(['Category', 'AssessmentType'])['Score'].mean().unstack()

            if avg_scores.empty:
                 logger.warning("Skipping Pre/Post plot: Average scores calculation resulted in empty data.")
                 return None

            # Handle potential missing columns (if a category only has Pre or Post)
            if 'Pre' not in avg_scores.columns:
                avg_scores['Pre'] = 0 # Add column with zeros if missing
            if 'Post' not in avg_scores.columns:
                 avg_scores['Post'] = 0 # Add column with zeros if missing

            # Sort by Improvement (optional)
            avg_scores['Improvement'] = avg_scores['Post'] - avg_scores['Pre']
            avg_scores = avg_scores.sort_values('Improvement', ascending=False)

            categories = avg_scores.index
            pre_scores = avg_scores['Pre'].values
            post_scores = avg_scores['Post'].values

            x = np.arange(len(categories))
            width = 0.35

            fig, ax = plt.subplots(figsize=(12, 7)) # Adjust figure size
            bars1 = ax.bar(x - width/2, pre_scores, width, label='Pre', color='navy')
            bars2 = ax.bar(x + width/2, post_scores, width, label='Post', color='lightgreen')

            # Annotate bars
            def autolabel(bars):
                for bar in bars:
                    height = bar.get_height()
                    ax.annotate(f'{height:.1f}%', # Format to 1 decimal place
                                xy=(bar.get_x() + bar.get_width() / 2, height),
                                xytext=(0, 3), # 3 points vertical offset
                                textcoords="offset points",
                                ha='center', va='bottom')

            autolabel(bars1)
            autolabel(bars2)

            # Labels and formatting
            ax.set_ylabel('Average Score (%)', fontsize=12) # Increase font size
            ax.set_title('Average Pre/Post Scores by Category', fontsize=14) # Increase font size
            ax.set_xticks(x)
            ax.set_xticklabels(categories, rotation=45, ha='right', fontsize=10) # Rotate and align labels
            ax.set_ylim(0, 100) # Assuming scores are percentages up to 100
            ax.legend()

            # Add grid for better readability
            ax.yaxis.grid(True, linestyle='--', alpha=0.7)

            plt.tight_layout() # Adjust layout to prevent labels overlapping

            plot_path = os.path.join(output_dir, 'pre_post_comparison.png')
            plt.savefig(plot_path)
            plt.close(fig) # Close the figure to free memory
            logger.info(f"Pre/Post comparison plot saved to {plot_path}")
            return plot_path
        except Exception as e:
            logger.error(f"Error creating Pre/Post comparison plot: {str(e)}")
            return None # Return None if plot generation fails

    def create_visualizations(self, output_dir: str) -> List[str]:
        """Create visualizations and save them to the output directory."""
        try:
            os.makedirs(output_dir, exist_ok=True)
            image_paths = []

            # Add the new Pre/Post comparison plot
            pre_post_plot_path = self._create_pre_post_comparison_plot(output_dir)
            if pre_post_plot_path:
                image_paths.append(pre_post_plot_path)

            # 1. Time Series Plot
            plt.figure(figsize=(12, 6))
            self.data.groupby('date')['score'].mean().plot()
            plt.title('Average Score Over Time')
            plt.xlabel('Date')
            plt.ylabel('Score')
            time_series_path = os.path.join(output_dir, 'time_series.png')
            plt.savefig(time_series_path)
            plt.close()
            image_paths.append(time_series_path)

            # 2. Distribution Plot
            plt.figure(figsize=(10, 6))
            sns.histplot(data=self.data, x='score', bins=30)
            plt.title('Score Distribution')
            plt.xlabel('Score')
            plt.ylabel('Count')
            dist_path = os.path.join(output_dir, 'distribution.png')
            plt.savefig(dist_path)
            plt.close()
            image_paths.append(dist_path)

            # 3. Category-wise Analysis (if category column exists)
            if 'category' in self.data.columns:
                plt.figure(figsize=(12, 6))
                sns.boxplot(data=self.data, x='category', y='score')
                plt.title('Score Distribution by Category')
                plt.xticks(rotation=45)
                category_path = os.path.join(output_dir, 'category_analysis.png')
                plt.savefig(category_path)
                plt.close()
                image_paths.append(category_path)

            return image_paths
        except Exception as e:
            logger.error(f"Error creating visualizations: {str(e)}")
            raise

    def generate_ppt_report(self, output_path: str, summary: Dict, image_paths: List[str]) -> str:
        """Generate PowerPoint report with analysis results."""
        try:
            prs = Presentation()

            # Title Slide
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = title_slide.shapes.title
            subtitle = title_slide.placeholders[1]
            title.text = "Daily Assessment Analysis Report"
            subtitle.text = f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"

            # Summary Slide
            summary_slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = summary_slide.shapes.title
            title.text = "Summary Statistics"
            
            # Add summary content
            content = summary_slide.placeholders[1]
            tf = content.text_frame
            tf.text = f"Total Records: {summary['total_records']}\n"
            tf.text += f"Date Range: {summary['date_range']['start']} to {summary['date_range']['end']}\n\n"
            
            for metric, values in summary['metrics'].items():
                tf.text += f"{metric}:\n"
                tf.text += f"  Mean: {values['mean']:.2f}\n"
                tf.text += f"  Median: {values['median']:.2f}\n"
                tf.text += f"  Std Dev: {values['std']:.2f}\n"

            # Visualization Slides
            for img_path in image_paths:
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                title = slide.shapes.title
                title.text = os.path.basename(img_path).replace('.png', '').replace('_', ' ').title()
                
                # Add image
                left = Inches(1)
                top = Inches(2)
                pic = slide.shapes.add_picture(img_path, left, top, width=Inches(8))

            # Save the presentation
            prs.save(output_path)
            logger.info(f"PowerPoint report generated: {output_path}")
            return output_path
        except Exception as e:
            logger.error(f"Error generating PowerPoint report: {str(e)}")
            raise

    def generate_report(self, report_id: str, filters: Optional[Dict[str, List[str]]] = None) -> str:
        """Generate complete analysis report."""
        try:
            # Apply filters if provided
            if filters:
                self.apply_filters(filters)

            # Create directory for plots (optional, can be deleted after report generation)
            plots_output_dir = os.path.join('reports', 'daily_assessment', 'temp_plots', datetime.now().strftime('%Y%m%d_%H%M%S'))
            os.makedirs(plots_output_dir, exist_ok=True)

            # Generate summary statistics
            summary = self.generate_summary_statistics()

            # Create visualizations
            image_paths = self.create_visualizations(plots_output_dir)

            # Create the final report directory if it doesn't exist
            final_report_dir = os.path.join('reports', 'daily_assessment')
            os.makedirs(final_report_dir, exist_ok=True)

            # Generate PowerPoint report with the report_id as the filename
            ppt_path = os.path.join(final_report_dir, f'{report_id}.pptx')
            self.generate_ppt_report(ppt_path, summary, image_paths)

            # Optional: Clean up temporary plot directory
            # shutil.rmtree(plots_output_dir)

            return ppt_path
        except Exception as e:
            logger.error(f"Error generating report: {str(e)}")
            raise

# Pydantic model for creating a new history entry
class HistoryCreate(BaseModel):
    user_id: str
    activity_type: str
    file_id: Optional[str] = None # Allow None for cases without a specific file
    file_name: Optional[str] = None # Allow None
    report_id: Optional[str] = None # Allow None
    filters: Optional[Dict] = None # Allow None
